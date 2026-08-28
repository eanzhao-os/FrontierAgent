"""Modal screens: tool-call approval and session resume.

These replace the line-mode stdin prompts. ``ApprovalScreen`` returns the same
:class:`Decision` the line-mode ``Approver`` produced, so the observer's
redirect / remember / reject branches are unchanged.
"""

from __future__ import annotations

import json
from dataclasses import dataclass
from pathlib import Path
from typing import Any, ClassVar, cast

from rich.console import RenderableType
from rich.markdown import Markdown
from rich.syntax import Syntax
from rich.table import Table
from rich.text import Text
from textual.app import ComposeResult
from textual.binding import Binding, BindingType
from textual.containers import Horizontal, Vertical, VerticalScroll
from textual.screen import ModalScreen
from textual.widgets import Button, Input, OptionList, Static, Tab, Tabs
from textual.widgets.option_list import Option

from apodex.observers import Decision, _target_suffix
from apodex.render import diff_to_text
from apodex.tui.themes import GLYPHS, active_theme, rich_style
from apodex.usage import Usage

_CODE_LEXERS = {
    ".py": "python", ".pyi": "python", ".ts": "typescript",
    ".tsx": "tsx", ".mts": "typescript", ".cts": "typescript",
    ".js": "javascript", ".jsx": "jsx", ".mjs": "javascript",
    ".cjs": "javascript", ".cs": "csharp",
    ".c": "c", ".h": "c", ".cpp": "cpp", ".cc": "cpp",
    ".hpp": "cpp", ".rs": "rust", ".go": "go", ".java": "java",
    ".kt": "kotlin", ".swift": "swift", ".rb": "ruby", ".php": "php",
    ".sh": "bash", ".bash": "bash", ".zsh": "bash", ".sql": "sql",
    ".json": "json", ".yaml": "yaml", ".yml": "yaml", ".toml": "toml",
    ".xml": "xml", ".html": "html", ".htm": "html", ".css": "css",
    ".scss": "scss", ".vue": "vue", ".svelte": "svelte",
}
_DOCUMENT_SUFFIXES = {
    ".md", ".markdown", ".txt", ".log", ".rst", ".csv", ".tsv", ".ini", ".cfg",
}
_BINARY_PREVIEWABLE_SUFFIXES = {
    ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt",
    ".ipynb", ".pdb", ".cif", ".ent", ".stl", ".obj", ".gltf", ".glb",
    ".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif",
    ".zip", ".tar", ".gz", ".tgz",
}
PREVIEWABLE_SUFFIXES = frozenset(_CODE_LEXERS) | _DOCUMENT_SUFFIXES | _BINARY_PREVIEWABLE_SUFFIXES
BINARY_PREVIEWABLE_SUFFIXES = _BINARY_PREVIEWABLE_SUFFIXES
_MAX_PREVIEW_BYTES = 500_000


class _Themed:
    """Resolve semantic roles for a modal from the app's active theme.

    Modal text is built as Rich ``Text``, so it does not pick up Textual's CSS
    variables — without this the modals kept Rich's generic ``red`` / ``cyan`` /
    ``dim`` and clashed with whichever theme was selected behind them.
    """

    def _style(self, role: str, *, bold: bool = False) -> str:
        return rich_style(active_theme(self), role, bold=bold)


@dataclass
class ApprovalOutcome:
    """What the approval modal returns.

    Wraps the engine-facing :class:`Decision` and adds ``all_session`` and
    ``auto_for_me`` — TUI/session concerns for auto-approval modes.
    """

    decision: Decision
    all_session: bool = False
    auto_for_me: bool = False


class ApprovalScreen(_Themed, ModalScreen[ApprovalOutcome]):
    """Ask the user to approve one tool call. Returns a :class:`Decision`.

    The proposed change (a unified diff, or the bash command) is shown *inside*
    the modal so the user can see exactly what they're approving. Normal calls
    are chosen either with arrow-keys + Enter on the option list, or with the
    single-key shortcuts y/n/m/a/A/e. ``e`` (or any dangerous call) reveals an
    input for a typed redirect; a *dangerous* call requires typing ``yes``.

    All text is built as Rich ``Text`` (never markup strings) so tool names,
    reasons, and the hotkey hints can't be mis-parsed as markup tags.
    """

    # We focus the option list (normal) or the input (dangerous) ourselves in
    # ``on_mount``; don't let Textual auto-focus something else.
    AUTO_FOCUS = None

    DEFAULT_CSS = """
    ApprovalScreen { align: center middle; }
    #ap-box {
        width: 80%; max-width: 100; height: auto; max-height: 90%; padding: 1 2;
        border: round $border; background: $surface; color: $text;
    }
    #ap-msg { height: auto; margin-bottom: 1; }
    #ap-preview {
        height: auto; max-height: 18; border: round $border-blurred; padding: 0 1;
        scrollbar-size-vertical: 1;
        margin-bottom: 1;
    }
    #ap-opts { height: auto; margin-bottom: 1; }
    #ap-hint { height: auto; }
    #ap-input { margin-top: 1; }
    """

    BINDINGS: ClassVar[list[BindingType]] = [
        ("y", "yes", "yes"),
        ("n", "no", "no"),
        ("m", "autome", "auto for me"),
        ("a", "all", "all"),
        ("A", "always", "always allow"),
        ("e", "redirect", "redirect"),
        ("ctrl+u", "preview_up", "scroll preview up"),
        ("ctrl+d", "preview_down", "scroll preview down"),
        ("escape", "no", "no"),
    ]

    # option id → the action that decides it
    _OPTIONS: ClassVar[Any] = [
        ("yes", "Yes — run it", "y"),
        ("no", "No — reject", "n"),
        ("autome", "Auto for me — auto-approve bash (docker/trusted env)", "m"),
        ("all", "Allow all for this session", "a"),
        ("always", "Always allow this command", "A"),
        ("redirect", "Redirect — tell it what to do instead", "e"),
    ]

    def __init__(
        self, tool_name: str, reason: str, dangerous: str = "",
        *, target: str = "", preview: str = "", preview_kind: str = "",
    ) -> None:
        super().__init__()
        self.tool_name = tool_name
        self.target = target
        self.reason = reason
        self.dangerous = dangerous
        self.preview = preview
        self.preview_kind = preview_kind

    def compose(self) -> ComposeResult:
        with Vertical(id="ap-box"):
            yield Static(self._header(), id="ap-msg")
            if self.preview:
                yield Static(self._preview_summary(), id="ap-summary")
                with VerticalScroll(id="ap-preview"):
                    yield Static(self._preview_renderable(), id="ap-preview-body")
            if not self.dangerous:
                yield OptionList(
                    *(Option(label, id=oid) for oid, label, _key in self._OPTIONS),
                    id="ap-opts",
                )
                yield Static(self._hint(), id="ap-hint")
            yield Input(id="ap-input")

    def on_mount(self) -> None:
        inp = self.query_one("#ap-input", Input)
        if self.dangerous:
            inp.placeholder = "type 'yes' to confirm, or an instruction to redirect"
            inp.focus()
        else:
            inp.display = False
            inp.can_focus = False  # keep single keys (y/n/a/A/e) at the screen level
            # Focus the option list so arrow-keys + Enter navigate. Its default
            # bindings are only arrows/enter/home/end/page — the letter shortcuts
            # bubble up to this screen's BINDINGS.
            options = self.query_one("#ap-opts", OptionList)
            # Approval should require an affirmative movement. Starting on No
            # prevents an accidental Enter from executing a write.
            options.highlighted = 1
            options.focus()

    def _header(self) -> Text:
        text = self._style("text")
        muted = self._style("muted")
        t = Text()
        if self.dangerous:
            t.append(f"{GLYPHS['danger']} DANGEROUS", style=self._style("err", bold=True))
            t.append(f" — {self.dangerous}\n", style=text)
            t.append("Confirm ", style=text)
            t.append(self.tool_name, style=self._style("tool", bold=True))
            t.append(_target_suffix(self.tool_name, self.target), style=text)
            t.append("?  ", style=text)
            t.append("type ", style=muted)
            t.append("yes", style=self._style("err", bold=True))
            t.append(" to confirm, or an instruction to redirect", style=muted)
            return t
        t.append("Approve ", style=text)
        t.append(self.tool_name, style=self._style("tool", bold=True))
        t.append(_target_suffix(self.tool_name, self.target), style=text)
        if self.reason:
            t.append(f"  ({self.reason})", style=muted)
        t.append("?", style=text)
        return t

    def _hint(self) -> Text:
        muted = self._style("muted")
        t = Text("↑↓ + Enter to choose  ·  shortcuts: ", style=muted)
        parts = [("y", "Yes"), ("n", "No"), ("a", "all"),
                 ("A", "Always"), ("e", "redirect")]
        for i, (key, lbl) in enumerate(parts):
            if i:
                t.append(" · ", style=self._style("subtle"))
            t.append(key, style=self._style("accent", bold=True))
            t.append(f" {lbl}", style=muted)
        return t

    def _preview_renderable(self) -> Text:
        theme = active_theme(self)
        if self.preview_kind == "diff":
            return diff_to_text(self.preview, theme=theme)
        t = Text()
        t.append("$ ", style=self._style("accent", bold=True))
        t.append(self.preview.strip(), style=self._style("text"))
        return t

    def _preview_summary(self) -> Text:
        """Compact scope/risk cue shown before the scrollable preview."""
        lines = self.preview.splitlines()
        if self.preview_kind == "diff":
            additions = sum(
                line.startswith("+") and not line.startswith("+++") for line in lines
            )
            deletions = sum(
                line.startswith("-") and not line.startswith("---") for line in lines
            )
            large = additions + deletions > 20
            return Text(
                f"Review change · {len(lines)} lines · +{additions} / -{deletions}"
                " · Ctrl-U/Ctrl-D scroll",
                style=self._style("warn", bold=True) if large else self._style("muted"),
            )
        return Text(
            f"Review command · {len(self.preview):,} characters · Ctrl-U/Ctrl-D scroll",
            style=self._style("warn", bold=True) if self.dangerous
            else self._style("muted"),
        )

    # ── option-list navigation (arrow-keys + Enter) ───────────────────────
    def on_option_list_option_selected(self, event: OptionList.OptionSelected) -> None:
        action = {
            "yes": self.action_yes, "no": self.action_no, "autome": self.action_autome,
            "all": self.action_all, "always": self.action_always,
            "redirect": self.action_redirect,
        }.get(event.option.id or "")
        if action is not None:
            action()

    # ── single-key actions (normal calls) ─────────────────────────────────
    def action_yes(self) -> None:
        if not self.dangerous:
            self.dismiss(ApprovalOutcome(Decision(True)))

    def action_no(self) -> None:
        self.dismiss(ApprovalOutcome(Decision(False)))

    def action_autome(self) -> None:
        if not self.dangerous:
            self.dismiss(ApprovalOutcome(Decision(True), auto_for_me=True))

    def action_all(self) -> None:
        if not self.dangerous:
            self.dismiss(ApprovalOutcome(Decision(True), all_session=True))

    def action_always(self) -> None:
        if not self.dangerous:
            self.dismiss(ApprovalOutcome(Decision(True, remember=True)))

    def action_redirect(self) -> None:
        if self.dangerous:
            return
        inp = self.query_one("#ap-input", Input)
        inp.placeholder = "tell the agent what to do instead"
        inp.display = True
        inp.can_focus = True
        inp.focus()

    def action_preview_up(self) -> None:
        if self.preview:
            self.query_one("#ap-preview", VerticalScroll).scroll_page_up(animate=False)

    def action_preview_down(self) -> None:
        if self.preview:
            self.query_one("#ap-preview", VerticalScroll).scroll_page_down(animate=False)

    # ── typed input (redirect / dangerous confirm) ────────────────────────
    def on_input_submitted(self, event: Input.Submitted) -> None:
        val = event.value.strip()
        if self.dangerous:
            low = val.lower()
            if low == "yes":
                self.dismiss(ApprovalOutcome(Decision(True)))
            elif low in ("", "y", "n", "no"):  # 'y' alone is intentionally NOT enough
                self.dismiss(ApprovalOutcome(Decision(False)))
            else:
                self.dismiss(ApprovalOutcome(Decision(False, feedback=val)))
        else:
            self.dismiss(ApprovalOutcome(
                Decision(False, feedback=val) if val else Decision(False)))


class ResumeScreen(ModalScreen[str]):
    """Pick a saved session to resume. Dismisses with the session id, or None."""

    DEFAULT_CSS = """
    ResumeScreen { align: center middle; }
    #rs-box {
        width: 80%; max-width: 100; height: auto; max-height: 80%; padding: 1 2;
        border: round $border; background: $surface; color: $text;
    }
    #rs-title { margin-bottom: 1; }
    """

    BINDINGS: ClassVar[list[BindingType]] = [("escape", "cancel", "cancel")]

    def compose(self) -> ComposeResult:
        with Vertical(id="rs-box"):
            yield Static("[bold]Resume a session[/]  ·  Esc to cancel", id="rs-title")
            yield OptionList(*self._options(), id="rs-list")

    @staticmethod
    def _options() -> list[Option]:
        from apodex.session_state import list_saved_sessions

        opts: list[Option] = []
        for saved in list_saved_sessions()[:20]:
            sid = saved["session_id"]
            name = str(saved.get("name") or "").strip()
            meta = (
                (f"{name} · " if name else "")
                + f"{saved['mode']}, {saved['message_count']} msgs"
            )
            opts.append(Option(f"{sid}  ({meta})", id=sid))
        if not opts:
            opts.append(Option("no saved sessions yet", disabled=True))
        return opts

    def on_option_list_option_selected(self, event: OptionList.OptionSelected) -> None:
        self.dismiss(event.option.id)

    def action_cancel(self) -> None:
        self.dismiss(None)


class ModelScreen(ModalScreen[str]):
    """Pick a model with arrow-keys + Enter. Dismisses with the model id, or None.

    Populated from the active profile's model list (``session.models``); the
    current model is pre-highlighted.
    """

    DEFAULT_CSS = """
    ModelScreen { align: center middle; }
    #ml-box {
        width: 80%; max-width: 100; height: auto; max-height: 80%; padding: 1 2;
        border: round $border; background: $surface; color: $text;
    }
    #ml-title { margin-bottom: 1; }
    """

    BINDINGS: ClassVar[list[BindingType]] = [("escape", "cancel", "cancel")]

    def __init__(self, models: list[str], current: str = "") -> None:
        super().__init__()
        self.models = models
        self.current = current

    def compose(self) -> ComposeResult:
        with Vertical(id="ml-box"):
            yield Static("[bold]Select a model[/]  ·  ↑↓ + Enter · Esc to cancel", id="ml-title")
            yield OptionList(*self._options(), id="ml-list")

    def _options(self) -> list[Option]:
        opts: list[Option] = []
        for m in self.models:
            label = m + ("  (current)" if m == self.current else "")
            opts.append(Option(label, id=m))
        if not opts:
            opts.append(Option("no models in this profile", disabled=True))
        return opts

    def on_mount(self) -> None:
        ol = self.query_one("#ml-list", OptionList)
        if self.current in self.models:
            ol.highlighted = self.models.index(self.current)
        ol.focus()

    def on_option_list_option_selected(self, event: OptionList.OptionSelected) -> None:
        self.dismiss(event.option.id)

    def action_cancel(self) -> None:
        self.dismiss(None)


class SelectionScreen(ModalScreen[str]):
    """A small mouse- and keyboard-selectable settings picker."""

    DEFAULT_CSS = """
    SelectionScreen { align: center middle; }
    #select-box {
        width: 80%; max-width: 76; height: auto; max-height: 80%; padding: 1 2;
        border: round $border; background: $surface; color: $text;
    }
    #select-title { margin-bottom: 1; }
    """
    BINDINGS: ClassVar[list[BindingType]] = [("escape", "cancel", "cancel")]

    def __init__(
        self, title: str, options: tuple[tuple[str, str], ...], current: str = "",
    ) -> None:
        super().__init__()
        self.title = title
        self.options = options
        self.current = current

    def compose(self) -> ComposeResult:
        with Vertical(id="select-box"):
            yield Static(
                f"[bold]{self.title}[/]  ·  ↑↓ + Enter or click · Esc to cancel",
                id="select-title",
            )
            yield OptionList(
                *(
                    Option(label + ("  (current)" if key == self.current else ""), id=key)
                    for key, label in self.options
                ),
                id="select-list",
            )

    def on_mount(self) -> None:
        options = self.query_one("#select-list", OptionList)
        keys = [key for key, _label in self.options]
        if self.current in keys:
            options.highlighted = keys.index(self.current)
        options.focus()

    def on_option_list_option_selected(self, event: OptionList.OptionSelected) -> None:
        self.dismiss(event.option.id)

    def action_cancel(self) -> None:
        self.dismiss(None)


class ThemeScreen(SelectionScreen):
    """Choose a registered TUI theme."""

    def __init__(self, themes: tuple[str, ...], current: str) -> None:
        super().__init__("Select theme", tuple((theme, theme) for theme in themes), current)


class WorkflowScreen(SelectionScreen):
    """Choose one of the native workflow profiles."""

    def __init__(self, current: str) -> None:
        super().__init__(
            "Select workflow",
            (
                ("react", "react  — single-agent Stateful ReAct"),
                ("agent_team", "agent_team  — coordinator with sub-agents"),
            ),
            current,
        )


@dataclass(frozen=True)
class SettingsOutcome:
    """The settings staged in :class:`SettingsScreen` when Apply is pressed."""

    theme: str
    workflow: str
    plan_mode: bool
    verbose: bool
    auto_approve: bool
    auto_for_me: bool = False
    resume_session_id: str = ""


class SettingsScreen(ModalScreen[SettingsOutcome]):
    """Unified, mouse-friendly settings menu for appearance and workflow."""

    AUTO_FOCUS = None
    DEFAULT_CSS = """
    SettingsScreen { align: center middle; }
    #settings-box {
        width: 86%; max-width: 104; height: 82%; max-height: 42; padding: 1 2;
        border: round $border; background: $surface; color: $text;
    }
    #settings-title { height: 2; color: $text; text-style: bold; }
    #settings-tabs { height: 3; margin-bottom: 1; }
    #settings-themes, #settings-workflows, #settings-behavior,
    #settings-permissions, #settings-sessions {
        height: 1fr; background: $surface; color: $text;
        scrollbar-size-vertical: 1;
    }
    #settings-hint { height: 2; color: $text-muted; padding-top: 1; }
    #settings-actions { height: 3; align-horizontal: center; }
    #settings-actions Button { min-width: 14; margin: 0 1; }
    """
    BINDINGS: ClassVar[list[BindingType]] = [
        Binding("left", "previous_section", "previous section", priority=True),
        Binding("right", "next_section", "next section", priority=True),
        Binding("space", "toggle_setting", "toggle", priority=True),
        Binding("enter", "apply", "apply", priority=True),
        Binding("escape", "cancel", "close", priority=True),
    ]

    _WORKFLOWS = (
        ("react", "react  — single-agent Stateful ReAct"),
        ("agent_team", "agent_team  — coordinator with sub-agents"),
    )

    def __init__(
        self,
        themes: tuple[str, ...],
        current_theme: str,
        current_workflow: str,
        *,
        plan_mode: bool = False,
        verbose: bool = True,
        auto_approve: bool = False,
        auto_for_me: bool = False,
        permission_allow: tuple[str, ...] = (),
        permission_deny: tuple[str, ...] = (),
        sessions: tuple[tuple[str, str], ...] = (),
        current_session: str = "",
    ) -> None:
        super().__init__()
        self.themes = themes
        self.selected_theme = current_theme
        self.selected_workflow = current_workflow
        self.plan_mode = plan_mode
        self.verbose = verbose
        self.auto_approve = auto_approve
        self.auto_for_me = auto_for_me
        self.permission_allow = permission_allow
        self.permission_deny = permission_deny
        self.sessions = sessions
        self.current_session = current_session
        self.selected_session = ""
        self._section = "theme"

    def compose(self) -> ComposeResult:
        with Vertical(id="settings-box"):
            yield Static("Settings", id="settings-title")
            yield Tabs(
                Tab("Theme", id="settings-tab-theme"),
                Tab("Workflow", id="settings-tab-workflow"),
                Tab("Behavior", id="settings-tab-behavior"),
                Tab("Permissions", id="settings-tab-permissions"),
                Tab("Sessions", id="settings-tab-sessions"),
                active="settings-tab-theme",
                id="settings-tabs",
            )
            yield OptionList(
                *(Option(self._theme_label(theme), id=theme) for theme in self.themes),
                id="settings-themes",
            )
            yield OptionList(
                *(
                    Option(self._workflow_label(key, label), id=key)
                    for key, label in self._WORKFLOWS
                ),
                id="settings-workflows",
            )
            yield OptionList(
                Option(self._toggle_label("Plan mode", self.plan_mode), id="behavior-plan"),
                Option(
                    self._toggle_label("Show full thinking", self.verbose),
                    id="behavior-verbose",
                ),
                id="settings-behavior",
            )
            yield OptionList(*self._permission_options(), id="settings-permissions")
            yield OptionList(*self._session_options(), id="settings-sessions")
            yield Static(
                "←→ section  ·  ↑↓ select  ·  Space toggle  ·  Enter apply  ·  Esc close",
                id="settings-hint",
            )
            with Horizontal(id="settings-actions"):
                yield Button("↵  Apply", id="settings-apply", variant="primary")
                yield Button("Esc  Close", id="settings-close")

    def on_mount(self) -> None:
        themes = self.query_one("#settings-themes", OptionList)
        workflows = self.query_one("#settings-workflows", OptionList)
        behavior = self.query_one("#settings-behavior", OptionList)
        permissions = self.query_one("#settings-permissions", OptionList)
        sessions = self.query_one("#settings-sessions", OptionList)
        for options in (workflows, behavior, permissions, sessions):
            options.display = False
        if self.selected_theme in self.themes:
            themes.highlighted = self.themes.index(self.selected_theme)
        workflow_keys = [key for key, _label in self._WORKFLOWS]
        if self.selected_workflow in workflow_keys:
            workflows.highlighted = workflow_keys.index(self.selected_workflow)
        themes.focus()

    def _theme_label(self, theme: str) -> str:
        return f"{theme}{'  ✓' if theme == self.selected_theme else ''}"

    def _workflow_label(self, key: str, label: str) -> str:
        return f"{label}{'  ✓' if key == self.selected_workflow else ''}"

    @staticmethod
    def _toggle_label(label: str, enabled: bool, *, warning: bool = False) -> str:
        value = "ON" if enabled else "OFF"
        suffix = "  ⚠ skips confirmation prompts" if warning and enabled else ""
        return f"{'✓' if enabled else '○'}  {label:<28} {value}{suffix}"

    def _permission_options(self) -> list[Option]:
        options = [
            Option(
                self._toggle_label(
                    "Bypass permissions (Auto-approve all)", self.auto_approve, warning=True,
                ),
                id="permissions-auto",
            ),
            Option(
                self._toggle_label(
                    "Auto for me (Auto-approve Bash & internal)", self.auto_for_me,
                ),
                id="permissions-autome",
            ),
            Option(
                f"Saved rules · {len(self.permission_allow)} allow / "
                f"{len(self.permission_deny)} deny",
                disabled=True,
            ),
        ]
        for rule in self.permission_allow:
            options.append(Option(f"  allow  {rule}", disabled=True))
        for rule in self.permission_deny:
            options.append(Option(f"  deny   {rule}", disabled=True))
        if not self.permission_allow and not self.permission_deny:
            options.append(Option("  No saved permission rules", disabled=True))
        return options

    def _session_options(self) -> list[Option]:
        options: list[Option] = []
        if self.current_session:
            options.append(Option(
                f"Current · {self.current_session}  ✓", id=self.current_session,
            ))
        options.extend(Option(label, id=session_id) for session_id, label in self.sessions)
        if not options:
            options.append(Option("No saved sessions yet", disabled=True))
        elif not self.sessions:
            options.append(Option("No other saved sessions", disabled=True))
        return options

    def _show_section(self, section: str) -> None:
        self._section = section
        sections = {
            "theme": "#settings-themes",
            "workflow": "#settings-workflows",
            "behavior": "#settings-behavior",
            "permissions": "#settings-permissions",
            "sessions": "#settings-sessions",
        }
        for name, selector in sections.items():
            self.query_one(selector, OptionList).display = name == section
        self.query_one("#settings-tabs", Tabs).active = f"settings-tab-{section}"
        self.query_one(sections[section], OptionList).focus()

    def on_tabs_tab_activated(self, event: Tabs.TabActivated) -> None:
        tab_id = event.tab.id or "settings-tab-theme"
        self._show_section(tab_id.removeprefix("settings-tab-"))

    def on_option_list_option_selected(self, event: OptionList.OptionSelected) -> None:
        """Stage mouse/keyboard selections; Apply commits both together."""
        selected = event.option.id
        if selected is None:
            return
        if event.option_list.id == "settings-themes":
            self.selected_theme = selected
            for theme in self.themes:
                event.option_list.replace_option_prompt(theme, self._theme_label(theme))
        elif event.option_list.id == "settings-workflows":
            self.selected_workflow = selected
            for key, label in self._WORKFLOWS:
                event.option_list.replace_option_prompt(
                    key, self._workflow_label(key, label),
                )
        else:
            self._activate_option(event.option_list, event.option)

    def _activate_option(self, options: OptionList, option: Option) -> None:
        if option.id == "behavior-plan":
            self.plan_mode = not self.plan_mode
            options.replace_option_prompt(
                option.id, self._toggle_label("Plan mode", self.plan_mode),
            )
        elif option.id == "behavior-verbose":
            self.verbose = not self.verbose
            options.replace_option_prompt(
                option.id, self._toggle_label("Show full thinking", self.verbose),
            )
        elif option.id == "permissions-auto":
            self.auto_approve = not self.auto_approve
            options.replace_option_prompt(
                option.id,
                self._toggle_label(
                    "Bypass permissions (Auto-approve all)", self.auto_approve, warning=True,
                ),
            )
        elif option.id == "permissions-autome":
            self.auto_for_me = not self.auto_for_me
            options.replace_option_prompt(
                option.id,
                self._toggle_label(
                    "Auto for me (Auto-approve Bash & internal)", self.auto_for_me,
                ),
            )
        elif options.id == "settings-sessions":
            self.selected_session = (
                "" if option.id == self.current_session else (option.id or "")
            )

    def _stage_highlighted(self) -> None:
        selectors = {
            "theme": "#settings-themes",
            "workflow": "#settings-workflows",
            "behavior": "#settings-behavior",
            "permissions": "#settings-permissions",
            "sessions": "#settings-sessions",
        }
        options = self.query_one(selectors[self._section], OptionList)
        option = options.highlighted_option
        if option is not None and option.id is not None:
            if self._section == "theme":
                self.selected_theme = option.id
            elif self._section == "workflow":
                self.selected_workflow = option.id
            elif self._section == "sessions":
                self.selected_session = (
                    "" if option.id == self.current_session else option.id
                )

    def action_previous_section(self) -> None:
        sections = ("theme", "workflow", "behavior", "permissions", "sessions")
        self._show_section(sections[(sections.index(self._section) - 1) % len(sections)])

    def action_next_section(self) -> None:
        sections = ("theme", "workflow", "behavior", "permissions", "sessions")
        self._show_section(sections[(sections.index(self._section) + 1) % len(sections)])

    def action_toggle_setting(self) -> None:
        if self._section not in ("behavior", "permissions"):
            return
        options = self.query_one(
            "#settings-behavior" if self._section == "behavior"
            else "#settings-permissions",
            OptionList,
        )
        option = options.highlighted_option
        if option is not None:
            self._activate_option(options, option)

    def action_apply(self) -> None:
        self._stage_highlighted()
        self.dismiss(SettingsOutcome(
            self.selected_theme,
            self.selected_workflow,
            self.plan_mode,
            self.verbose,
            self.auto_approve,
            self.auto_for_me,
            self.selected_session,
        ))

    def action_cancel(self) -> None:
        self.dismiss(None)

    def on_button_pressed(self, event: Button.Pressed) -> None:
        if event.button.id == "settings-apply":
            self.action_apply()
        elif event.button.id == "settings-close":
            self.action_cancel()


class CommandScreen(ModalScreen[str]):
    """Keyboard-first slash-command palette."""

    DEFAULT_CSS = """
    CommandScreen { align: center middle; }
    #cmd-box {
        width: 80%; max-width: 90; height: auto; max-height: 85%; padding: 1 2;
        border: round $border; background: $surface; color: $text;
    }
    #cmd-title { margin-bottom: 1; }
    """
    BINDINGS: ClassVar[list[BindingType]] = [("escape", "cancel", "cancel")]

    def __init__(self, commands: tuple[tuple[str, str], ...]) -> None:
        super().__init__()
        self.commands = commands

    def compose(self) -> ComposeResult:
        with Vertical(id="cmd-box"):
            yield Static("[bold]Commands[/]  ·  ↑↓ + Enter · Esc to cancel", id="cmd-title")
            yield OptionList(
                *(
                    Option(f"{command:<10} {description}", id=command)
                    for command, description in self.commands
                ),
                id="cmd-list",
            )

    def on_mount(self) -> None:
        self.query_one("#cmd-list", OptionList).focus()

    def on_option_list_option_selected(self, event: OptionList.OptionSelected) -> None:
        self.dismiss(event.option.id)

    def action_cancel(self) -> None:
        self.dismiss(None)


class HelpScreen(_Themed, ModalScreen[None]):
    """Discoverable shortcut and interaction reference."""

    DEFAULT_CSS = """
    HelpScreen { align: center middle; }
    #help-box {
        width: 80%; max-width: 90; height: auto; max-height: 85%; padding: 1 2;
        border: round $border; background: $surface; color: $text;
    }
    #help-body { height: auto; max-height: 30; scrollbar-size-vertical: 1; }
    """
    BINDINGS: ClassVar[list[BindingType]] = [("escape", "close", "close"), ("f1", "close", "close")]

    def compose(self) -> ComposeResult:
        heading = self._style("accent", bold=True)
        key_style = self._style("user", bold=True)
        text = self._style("text")
        body = Text()
        body.append("Keyboard shortcuts\n", style=heading)
        for key, label in (
            ("F1", "this help"),
            ("F2", "theme and workflow settings"),
            ("Ctrl-P", "command palette"),
            ("Ctrl-B", "toggle sidebar"),
            ("Ctrl-Tab", "cycle sidebar tabs (Diff appears when files changed)"),
            ("Ctrl-O", "jump to Files; press again to return to Plan"),
            ("Space", "open the selected tool call or file"),
            ("Alt-J / K", "move through transcript blocks"),
            ("Alt-Enter", "expand or collapse the selected block"),
            ("Ctrl-G", "jump to the final report"),
            ("Ctrl-Y", "copy the final report"),
            ("Ctrl-V", "paste Finder files or an image from the macOS clipboard"),
            ("↑ / ↓", "prompt history"),
            ("Tab", "complete a slash command or @ file"),
            ("Ctrl-C", "interrupt current task; quit when idle"),
        ):
            body.append(f"  {key:<10}", style=key_style)
            body.append(label + "\n", style=text)
        body.append("\nWhile the agent works\n", style=heading)
        body.append("  Type a message to steer it at the next turn boundary.\n", style=text)
        body.append("  Slash commands wait until the current task is interrupted.\n", style=text)
        body.append("  Ctrl-P lists commands; /config shows safe local settings.\n", style=text)
        body.append("  /new saves this session; /context shows window usage.\n", style=text)
        body.append(
            "  /attach <path> adds read-only files; relative paths start at --cwd.\n",
            style=text,
        )
        body.append("  @ searches attached and workspace files under the current cwd.\n", style=text)
        body.append("  /paste reads Finder files or an image from the macOS clipboard.\n", style=text)
        body.append(
            "  Multiline/large text paste is folded in the prompt and sent in full.\n",
            style=text,
        )
        body.append("\nApproval\n", style=heading)
        body.append("  No is selected by default. y/n/a/A/e are shortcuts.\n", style=text)
        body.append("  Ctrl-U/Ctrl-D scroll a long command or diff preview.\n", style=text)
        body.append("\nEsc or F1 closes this help.", style=self._style("muted"))
        with Vertical(id="help-box"), VerticalScroll(id="help-body"):
            yield Static(body)

    def action_close(self) -> None:
        self.dismiss(None)


class ContextScreen(_Themed, ModalScreen[None]):
    """Current context-window usage with an estimated category distribution."""

    AUTO_FOCUS = None
    DEFAULT_CSS = """
    ContextScreen { align: center middle; }
    #context-box {
        width: 72; max-width: 92%; height: auto; max-height: 85%; padding: 1 2;
        border: round $border; background: $surface; color: $text;
    }
    #context-body { height: auto; }
    #context-hint { height: 1; margin-top: 1; color: $text-muted; }
    """
    BINDINGS: ClassVar[list[BindingType]] = [
        Binding("escape", "close", "Close", priority=True),
        Binding("space", "close", "Close", priority=True),
    ]

    def __init__(
        self, usage: Usage, *, window: int, output_reserve: int = 0,
        model: str = "",
    ) -> None:
        super().__init__()
        self.usage = usage
        self.window = max(0, int(window))
        self.output_reserve = max(0, int(output_reserve))
        self.model = model

    def compose(self) -> ComposeResult:
        body = Text()
        heading = self._style("accent", bold=True)
        text = self._style("text")
        muted = self._style("muted")
        body.append("Context usage", style=heading)
        if self.model:
            body.append(f"  ·  {self.model}", style=muted)
        body.append("\n\n")

        if self.window <= 0 or self.usage.last_input <= 0:
            limit = f" / {self.window:,}" if self.window else ""
            body.append(f"Current input unavailable{limit}\n", style=text)
            body.append(
                "Usage becomes available after the first model response.",
                style=muted,
            )
        else:
            used = min(self.usage.last_input, self.window)
            percent = min(100, round(used / self.window * 100, 1))
            marker = "≈" if self.usage.estimated else ""
            body.append(
                f"{marker}{used:,} / {self.window:,} tokens  ·  {percent:.1f}%\n",
                style=text,
            )
            self._append_bar(body, used)
            body.append("\n\n")
            breakdown = self.usage.breakdown
            if breakdown is not None:
                roles = ("muted", "accent", "tool", "warn")
                markers = ("◆", "◆", "◆", "◆")
                for (label, tokens), role, glyph in zip(
                    breakdown.display_categories(), roles, markers, strict=True,
                ):
                    pct = tokens / self.window * 100
                    body.append(f"{glyph} ", style=self._style(role, bold=True))
                    body.append(f"{label:<25}", style=text)
                    body.append(f"{tokens:>10,}  {pct:>5.1f}%\n", style=muted)
            free = max(0, self.window - used)
            body.append("◇ ", style=self._style("subtle"))
            body.append(f"{'Free':<25}", style=text)
            body.append(
                f"{free:>10,}  {free / self.window * 100:>5.1f}%\n", style=muted,
            )
            if self.output_reserve:
                body.append(
                    f"\nOutput reserve  up to {self.output_reserve:,} tokens",
                    style=muted,
                )
            compact = (
                str(self.usage.compactions)
                if self.usage.compactions else "not triggered"
            )
            body.append(f"\nCompaction      {compact}", style=muted)
            body.append(
                f"\nSession total   {self.usage.total:,} tokens",
                style=muted,
            )
            if self.usage.cached:
                body.append(f"  ·  cache read {self.usage.cached:,}", style=muted)

        with Vertical(id="context-box"):
            yield Static(body, id="context-body")
            yield Static(
                "Estimated categories · provider-reported total · Esc to close",
                id="context-hint",
            )

    def _append_bar(self, body: Text, used: int) -> None:
        width = 48
        breakdown = self.usage.breakdown
        if breakdown is None:
            filled = min(width, round(used / self.window * width))
            body.append("█" * filled, style=self._style("accent", bold=True))
            body.append("░" * (width - filled), style=self._style("subtle"))
            return
        roles = ("muted", "accent", "tool", "warn")
        cells_left = width
        for (_label, tokens), role in zip(
            breakdown.display_categories(), roles, strict=True,
        ):
            cells = min(cells_left, round(tokens / self.window * width))
            if cells:
                body.append("█" * cells, style=self._style(role, bold=True))
                cells_left -= cells
        if cells_left:
            body.append("░" * cells_left, style=self._style("subtle"))

    def action_close(self) -> None:
        self.dismiss(None)


class FilePreviewScreen(_Themed, ModalScreen[None]):
    """Read-only preview for a text or source-code deliverable."""

    AUTO_FOCUS = None
    DEFAULT_CSS = """
    FilePreviewScreen { align: center middle; }
    #file-preview-box {
        width: 82%; max-width: 110; height: 76%; max-height: 34; padding: 0 1;
        border: round $border; background: $surface; color: $text;
    }
    #file-preview-title {
        height: 2; padding: 0 1; color: $primary; text-style: bold;
        border-bottom: solid $border;
    }
    #file-preview-scroll {
        height: 1fr; padding: 1; scrollbar-size-vertical: 1;
        background: $background;
    }
    #file-preview-content { height: auto; }
    #file-preview-hint { height: 1; color: $text-muted; padding: 0 1; }
    """
    BINDINGS: ClassVar[list[BindingType]] = [
        Binding("space", "close", "Close", priority=True),
        Binding("escape", "close", "Close", priority=True),
        Binding("ctrl+u", "page_up", "Page up", priority=True),
        Binding("ctrl+d", "page_down", "Page down", priority=True),
    ]

    def __init__(self, path: Path, *, label: str) -> None:
        super().__init__()
        self.path = path
        self.label = label

    def _fallback_metadata(self, file_type: str, pkg_name: str, install_cmd: str) -> Text:
        t = Text()
        t.append(f"📄 {file_type} · {self.path.name}\n", style=self._style("accent", bold=True))
        t.append("─" * 50 + "\n", style=self._style("muted"))
        try:
            size_mb = self.path.stat().st_size / (1024 * 1024)
            t.append(f"File size: {size_mb:.2f} MB ({self.path.stat().st_size:,} bytes)\n", style=self._style("text"))
        except OSError:
            pass
        t.append(f"Path: {self.path}\n\n", style=self._style("muted"))
        t.append(f"💡 Install optional package '{pkg_name}' for full content preview:\n", style=self._style("warning"))
        t.append(f"   {install_cmd}\n", style=self._style("text", bold=True))
        return t

    def _render_pdf(self) -> RenderableType:
        try:
            import pypdf
            reader = pypdf.PdfReader(str(self.path))
            num_pages = len(reader.pages)
            meta = reader.metadata or {}
            title = meta.get("/Title", "") or self.path.name
            author = meta.get("/Author", "")

            t = Text()
            t.append(f"📄 PDF Document · {num_pages} page(s)\n", style=self._style("accent", bold=True))
            if title and title != self.path.name:
                t.append(f"Title: {title}\n", style=self._style("text", bold=True))
            if author:
                t.append(f"Author: {author}\n", style=self._style("muted"))
            t.append("─" * 45 + "\n\n", style=self._style("muted"))

            max_pages = min(num_pages, 5)
            for i in range(max_pages):
                t.append(f"--- Page {i + 1} ---\n", style=self._style("warning", bold=True))
                page_text = reader.pages[i].extract_text() or ""
                if page_text.strip():
                    t.append(page_text.strip()[:2000] + "\n\n", style=self._style("text"))
                else:
                    t.append("[No extractable text on this page]\n\n", style=self._style("muted"))
            if num_pages > max_pages:
                t.append(f"… {num_pages - max_pages} more page(s) not shown.\n", style=self._style("muted"))
            return t
        except ImportError:
            return self._fallback_metadata("PDF Document (.pdf)", "pypdf", "pip install frontier-agent[document-readers]")
        except Exception as exc:
            return Text(f"Could not read PDF: {exc}", style=self._style("muted"))

    def _render_docx(self) -> RenderableType:
        try:
            import docx
            doc = docx.Document(str(self.path))
            lines = [
                f"# {self.path.name}\n",
                f"*Word Document ({len(doc.paragraphs)} paragraphs, {len(doc.tables)} tables)*\n",
                "---",
            ]
            for p in doc.paragraphs:
                if not p.text.strip():
                    continue
                style_name = (p.style.name or "").lower() if p.style else ""
                if "heading 1" in style_name:
                    lines.append(f"\n# {p.text}")
                elif "heading 2" in style_name:
                    lines.append(f"\n## {p.text}")
                elif "heading 3" in style_name:
                    lines.append(f"\n### {p.text}")
                else:
                    lines.append(p.text)
            return Markdown("\n\n".join(lines))
        except ImportError:
            return self._fallback_metadata("Word Document (.docx)", "python-docx", "pip install frontier-agent[document-readers]")
        except Exception as exc:
            return Text(f"Could not read Word document: {exc}", style=self._style("muted"))

    def _render_xlsx(self) -> RenderableType:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(filename=str(self.path), read_only=True, data_only=True)
            sheet_names = wb.sheetnames
            sheet_name = sheet_names[0] if sheet_names else "Sheet1"
            ws = wb[sheet_name]
            table = Table(title=f"Excel: {self.path.name} ({len(sheet_names)} sheet(s): {', '.join(sheet_names[:3])})", show_header=True, header_style="bold magenta")
            rows = list(ws.iter_rows(values_only=True))
            if rows:
                header = rows[0]
                for i, col in enumerate(header[:8]):
                    table.add_column(str(col if col is not None else f"Col {i+1}"))
                for row in rows[1:25]:
                    table.add_row(*[str(cell if cell is not None else "") for cell in row[:8]])
            wb.close()
            return table
        except ImportError:
            return self._fallback_metadata("Excel Spreadsheet (.xlsx)", "openpyxl", "pip install frontier-agent[sandbox]")
        except Exception as exc:
            return Text(f"Could not read Excel file: {exc}", style=self._style("muted"))

    def _render_pptx(self) -> RenderableType:
        try:
            import pptx
            prs = pptx.Presentation(str(self.path))
            lines = [
                f"# {self.path.name}\n",
                f"*PowerPoint Presentation ({len(prs.slides)} slides)*\n",
                "---",
            ]
            for idx, slide in enumerate(prs.slides, 1):
                slide_title = "Untitled Slide"
                if slide.shapes.title and slide.shapes.title.text:
                    slide_title = slide.shapes.title.text.strip()
                lines.append(f"\n### Slide {idx}: {slide_title}")
                body_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        # has_text_frame is python-pptx's own guard for this
                        # access; BaseShape does not declare text_frame.
                        for paragraph in shape.text_frame.paragraphs:  # pyright: ignore[reportAttributeAccessIssue]
                            if paragraph.text.strip():
                                body_texts.append(f"- {paragraph.text.strip()}")
                if body_texts:
                    lines.extend(body_texts)
            return Markdown("\n".join(lines))
        except ImportError:
            return self._fallback_metadata("PowerPoint Presentation (.pptx)", "python-pptx", "pip install frontier-agent[document-readers]")
        except Exception as exc:
            return Text(f"Could not read PowerPoint presentation: {exc}", style=self._style("muted"))

    def _render_ipynb(self) -> RenderableType:
        try:
            with self.path.open("r", encoding="utf-8", errors="replace") as f:
                nb = json.load(f)
            cells = nb.get("cells", [])
            t = Text()
            t.append(f"📓 Jupyter Notebook · {len(cells)} cell(s)\n", style=self._style("accent", bold=True))
            t.append("─" * 45 + "\n\n", style=self._style("muted"))
            for idx, cell in enumerate(cells[:20], 1):
                cell_type = cell.get("cell_type", "code")
                source = "".join(cell.get("source", []))
                if cell_type == "markdown":
                    t.append(f"─── [Cell {idx}: Markdown] ───\n", style=self._style("accent"))
                    t.append(source + "\n\n", style=self._style("text"))
                elif cell_type == "code":
                    t.append(f"─── [Cell {idx}: Code] ───\n", style=self._style("warning"))
                    t.append(source + "\n", style=self._style("text"))
                    outputs = cell.get("outputs", [])
                    for out in outputs[:3]:
                        out_type = out.get("output_type", "")
                        if out_type in ("stream", "execute_result"):
                            text_out = "".join(out.get("text", []) or out.get("data", {}).get("text/plain", []))
                            if text_out.strip():
                                t.append("Out: " + text_out.strip()[:400] + "\n", style=self._style("muted"))
                    t.append("\n")
            if len(cells) > 20:
                t.append(f"… {len(cells) - 20} more cell(s) omitted.\n", style=self._style("muted"))
            return t
        except Exception as exc:
            return Text(f"Could not read Jupyter notebook: {exc}", style=self._style("muted"))

    def _render_pdb(self) -> RenderableType:
        try:
            with self.path.open("r", encoding="utf-8", errors="replace") as f:
                lines = f.readlines(_MAX_PREVIEW_BYTES)
            header = ""
            title = ""
            expdta = ""
            resolution = ""
            chains = set()
            atom_count = 0
            hetatm_count = 0
            b_factors = []

            for line in lines:
                rec = line[:6].strip()
                if rec == "HEADER":
                    header = line[10:50].strip()
                elif rec == "TITLE":
                    title += line[10:70].strip() + " "
                elif rec == "EXPDTA":
                    expdta = line[10:70].strip()
                elif rec == "REMARK" and "RESOLUTION." in line:
                    resolution = line.strip()
                elif rec in ("ATOM", "HETATM"):
                    if rec == "ATOM":
                        atom_count += 1
                    else:
                        hetatm_count += 1
                    chain_id = line[21:22].strip()
                    if chain_id:
                        chains.add(chain_id)
                    try:
                        bf = float(line[60:66].strip())
                        b_factors.append(bf)
                    except ValueError:
                        pass

            t = Text()
            t.append(f"🧬 PDB Biological Structure · {self.path.name}\n", style=self._style("accent", bold=True))
            t.append("─" * 45 + "\n", style=self._style("muted"))
            if header:
                t.append(f"Header: {header}\n", style=self._style("text", bold=True))
            if title:
                t.append(f"Title: {title.strip()}\n", style=self._style("text"))
            if expdta:
                t.append(f"Method: {expdta}\n", style=self._style("muted"))
            if resolution:
                t.append(f"{resolution}\n", style=self._style("muted"))
            t.append(f"Chains: {', '.join(sorted(chains)) if chains else 'N/A'}\n", style=self._style("text"))
            t.append(f"Atoms: {atom_count:,} (ATOM) | {hetatm_count:,} (HETATM)\n", style=self._style("text"))
            if b_factors:
                avg_b = sum(b_factors) / len(b_factors)
                t.append(f"Avg B-Factor / pLDDT: {avg_b:.2f}\n", style=self._style("accent", bold=True))
            return t
        except Exception as exc:
            return Text(f"Could not read PDB file: {exc}", style=self._style("muted"))

    def _render_3d(self) -> RenderableType:
        suffix = self.path.suffix.lower()
        try:
            t = Text()
            t.append(f"🧊 3D Asset · {self.path.name} ({suffix.upper()})\n", style=self._style("accent", bold=True))
            t.append("─" * 45 + "\n", style=self._style("muted"))

            if suffix == ".obj":
                with self.path.open("r", encoding="utf-8", errors="replace") as f:
                    content = f.read()
                v_count = sum(1 for line in content.splitlines() if line.startswith("v "))
                f_count = sum(1 for line in content.splitlines() if line.startswith("f "))
                t.append(f"Format: Wavefront OBJ\nVertices: {v_count:,}\nFaces: {f_count:,}\n", style=self._style("text"))
            elif suffix == ".stl":
                with self.path.open("rb") as f:
                    header = f.read(80)
                    is_ascii = header.startswith(b"solid")
                if is_ascii:
                    with self.path.open("r", encoding="utf-8", errors="replace") as f:
                        facets = sum(1 for line in f if "facet normal" in line)
                    t.append(f"Format: ASCII STL\nFacets: {facets:,}\n", style=self._style("text"))
                else:
                    with self.path.open("rb") as f:
                        f.seek(80)
                        import struct
                        num_triangles = struct.unpack("<I", f.read(4))[0]
                    t.append(f"Format: Binary STL\nTriangles: {num_triangles:,}\n", style=self._style("text"))
            elif suffix in (".gltf", ".glb"):
                if suffix == ".gltf":
                    with self.path.open("r", encoding="utf-8", errors="replace") as f:
                        data = json.load(f)
                    asset = data.get("asset", {})
                    t.append(f"Format: glTF {asset.get('version', '2.0')}\n", style=self._style("text"))
                    t.append(f"Meshes: {len(data.get('meshes', []))}\nMaterials: {len(data.get('materials', []))}\nNodes: {len(data.get('nodes', []))}\n", style=self._style("text"))
                else:
                    t.append(f"Format: Binary glTF (GLB)\nSize: {self.path.stat().st_size:,} bytes\n", style=self._style("text"))
            return t
        except Exception as exc:
            return Text(f"Could not read 3D asset: {exc}", style=self._style("muted"))

    def _render_archive(self) -> RenderableType:
        suffix = self.path.suffix.lower()
        try:
            t = Text()
            t.append(f"📦 Archive · {self.path.name}\n", style=self._style("accent", bold=True))
            t.append("─" * 45 + "\n", style=self._style("muted"))
            if suffix == ".zip":
                import zipfile
                with zipfile.ZipFile(self.path, "r") as z:
                    infos = z.infolist()
                    t.append(f"Entries: {len(infos)}\n\n", style=self._style("text", bold=True))
                    for info in infos[:25]:
                        size_str = f"{info.file_size:,} B" if not info.is_dir() else "<DIR>"
                        t.append(f"{info.filename:<40} {size_str:>12}\n", style=self._style("text"))
                    if len(infos) > 25:
                        t.append(f"\n… {len(infos) - 25} more entries omitted.\n", style=self._style("muted"))
            elif suffix in (".tar", ".gz", ".tgz"):
                import tarfile
                mode = "r:gz" if suffix in (".gz", ".tgz") else "r"
                with tarfile.open(self.path, mode) as tar:
                    members = tar.getmembers()
                    t.append(f"Entries: {len(members)}\n\n", style=self._style("text", bold=True))
                    for m in members[:25]:
                        size_str = f"{m.size:,} B" if m.isfile() else "<DIR>"
                        t.append(f"{m.name:<40} {size_str:>12}\n", style=self._style("text"))
                    if len(members) > 25:
                        t.append(f"\n… {len(members) - 25} more entries omitted.\n", style=self._style("muted"))
            return t
        except Exception as exc:
            return Text(f"Could not read archive: {exc}", style=self._style("muted"))

    def _render_image(self) -> RenderableType:
        try:
            from PIL import Image
            with Image.open(self.path) as img:
                w, h = img.size
                mode = img.mode
                t = Text()
                t.append(f"🖼️ Image · {w}x{h} ({mode}) · {self.path.name}\n", style=self._style("accent", bold=True))
                t.append("─" * 45 + "\n", style=self._style("muted"))

                target_w = min(w, 60)
                aspect = h / w if w > 0 else 1.0
                target_h = max(1, int(target_w * aspect * 0.5))
                target_h = min(target_h, 20)

                resized = img.convert("RGB").resize((target_w, target_h * 2), Image.Resampling.BILINEAR)
                pixels = resized.load()
                if pixels is None:
                    raise ValueError("image pixel data unavailable")

                # convert("RGB") above guarantees 3-int pixels, but the pixel
                # accessor is typed for every mode, so each read is narrowed.
                for y in range(0, target_h * 2, 2):
                    for x in range(target_w):
                        r1, g1, b1 = cast("tuple[int, int, int]", pixels[x, y])
                        r2, g2, b2 = (
                            cast("tuple[int, int, int]", pixels[x, y + 1])
                            if (y + 1) < target_h * 2 else (0, 0, 0)
                        )
                        t.append("▀", style=f"rgb({r1},{g1},{b1}) on rgb({r2},{g2},{b2})")
                    t.append("\n")
                return t
        except ImportError:
            return self._fallback_metadata("Image File", "Pillow", "pip install Pillow")
        except Exception as exc:
            return Text(f"Could not load image: {exc}", style=self._style("muted"))

    def _render_csv(self) -> RenderableType:
        try:
            import csv
            with self.path.open("r", encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f)
                rows = list(reader)
            if not rows:
                return Text("Empty CSV file", style=self._style("muted"))
            table = Table(show_header=True, header_style="bold magenta")
            header = rows[0]
            for i, col in enumerate(header[:10]):
                table.add_column(str(col) if col else f"Col {i+1}")
            for row in rows[1:30]:
                table.add_row(*[str(cell) for cell in row[:10]])
            return table
        except Exception as exc:
            return Text(f"Could not parse CSV: {exc}", style=self._style("muted"))

    def _renderable(self) -> RenderableType:
        from apodex.preview import build_preview
        from apodex.tui.themes import palette

        prev = build_preview(self.path)
        family = prev["family"]
        text = prev["text"] or ""
        error = prev["error"]
        title = prev["title"] or self.path.name
        if error and not text:
            return Text(error, style=self._style("muted"))
        if family in {"markdown", "docx", "pptx"}:
            return Markdown(text)
        if family == "csv":
            rows = prev["metadata"].get("rows") or []
            if rows:
                table = Table(show_header=True, header_style="bold magenta")
                header = rows[0]
                for i, col in enumerate(header[:10]):
                    table.add_column(str(col) if col else f"Col {i + 1}")
                for row in rows[1:30]:
                    table.add_row(*[str(cell) for cell in row[:10]])
                return table
        if family == "xlsx":
            rows = prev["metadata"].get("rows") or []
            sheets = prev["metadata"].get("sheets") or []
            table = Table(
                title=f"Excel: {self.path.name} ({len(sheets)} sheet(s): {', '.join(sheets[:3])})",
                show_header=True,
                header_style="bold magenta",
            )
            if rows:
                header = rows[0]
                for i, col in enumerate(header[:8]):
                    table.add_column(str(col if col is not None else f"Col {i + 1}"))
                for row in rows[1:25]:
                    table.add_row(*[str(cell if cell is not None else "") for cell in row[:8]])
            return table
        lexer = _CODE_LEXERS.get(self.path.suffix.lower())
        if lexer and family == "text":
            theme = "ansi_dark" if palette(active_theme(self)).dark else "ansi_light"
            return Syntax(
                text, lexer, theme=theme, background_color="default",
                line_numbers=True, word_wrap=False,
            )
        body = text
        if title and title not in body:
            body = f"{title}\n{body}"
        return Text(body, style=self._style("text"))

    def compose(self) -> ComposeResult:
        with Vertical(id="file-preview-box"):
            yield Static(f"Preview · {self.label}", id="file-preview-title")
            with VerticalScroll(id="file-preview-scroll"):
                yield Static(self._renderable(), id="file-preview-content")
            yield Static("Space / Esc close · Ctrl-U / Ctrl-D page", id="file-preview-hint")

    def action_close(self) -> None:
        self.dismiss(None)

    def action_page_up(self) -> None:
        self.query_one("#file-preview-scroll", VerticalScroll).action_page_up()

    def action_page_down(self) -> None:
        self.query_one("#file-preview-scroll", VerticalScroll).action_page_down()


class ActivityDetailScreen(_Themed, ModalScreen[None]):
    """Compact detail view for one tool call in the Activity tab."""

    AUTO_FOCUS = None
    DEFAULT_CSS = """
    ActivityDetailScreen { align: center middle; }
    #activity-detail-box {
        width: 72%; max-width: 88; height: auto; max-height: 70%; padding: 1 2;
        border: round $border; background: $surface; color: $text;
    }
    #activity-detail-title {
        height: auto; color: $primary; text-style: bold; margin-bottom: 1;
    }
    #activity-detail-body {
        height: auto; max-height: 18; padding: 0 1;
        border-left: solid $border; scrollbar-size-vertical: 1;
    }
    #activity-detail-hint { height: 1; color: $text-muted; margin-top: 1; }
    """
    BINDINGS: ClassVar[list[BindingType]] = [
        Binding("space", "close", "Close", priority=True),
        Binding("escape", "close", "Close", priority=True),
    ]

    def __init__(
        self, *, name: str, call_id: str, state: str, duration: str, summary: str,
    ) -> None:
        super().__init__()
        self.tool_name = name
        self.call_id = call_id
        self.state = state
        self.duration = duration
        self.summary = summary

    def compose(self) -> ComposeResult:
        body = Text()
        label = self._style("muted")
        value = self._style("text")
        for key, content in (
            ("State", self.state), ("Duration", self.duration),
            ("Call ID", self.call_id), ("Details", self.summary or "(none)"),
        ):
            body.append(f"{key:<10}", style=label)
            body.append(content + "\n", style=value)
        with Vertical(id="activity-detail-box"):
            yield Static(f"Tool call · {self.tool_name}", id="activity-detail-title")
            with VerticalScroll(id="activity-detail-body"):
                yield Static(body)
            yield Static("Space / Esc close", id="activity-detail-hint")

    def action_close(self) -> None:
        self.dismiss(None)


__all__ = [
    "BINARY_PREVIEWABLE_SUFFIXES",
    "PREVIEWABLE_SUFFIXES",
    "ActivityDetailScreen",
    "ApprovalScreen",
    "CommandScreen",
    "FilePreviewScreen",
    "HelpScreen",
    "ModelScreen",
    "ResumeScreen",
    "SelectionScreen",
    "SettingsOutcome",
    "SettingsScreen",
    "ThemeScreen",
    "WorkflowScreen",
]
