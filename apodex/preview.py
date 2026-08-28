"""Backend-neutral file preview classification and structured content."""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any

_DEFAULT_MAX_BYTES = 500_000

_SUFFIX_FAMILY = {
    ".md": "markdown",
    ".markdown": "markdown",
    ".csv": "csv",
    ".pdf": "pdf",
    ".docx": "docx",
    ".doc": "docx",
    ".xlsx": "xlsx",
    ".xls": "xlsx",
    ".pptx": "pptx",
    ".ppt": "pptx",
    ".ipynb": "notebook",
    ".png": "image",
    ".jpg": "image",
    ".jpeg": "image",
    ".webp": "image",
    ".bmp": "image",
    ".gif": "image",
    ".zip": "archive",
    ".tar": "archive",
    ".gz": "archive",
    ".tgz": "archive",
    ".pdb": "pdb",
    ".cif": "pdb",
    ".ent": "pdb",
    ".stl": "model3d",
    ".obj": "model3d",
    ".gltf": "model3d",
    ".glb": "model3d",
    ".txt": "text",
    ".py": "text",
    ".js": "text",
    ".ts": "text",
    ".json": "text",
    ".toml": "text",
    ".yml": "text",
    ".yaml": "text",
    ".html": "text",
    ".css": "text",
    ".rs": "text",
    ".go": "text",
    ".java": "text",
    ".c": "text",
    ".h": "text",
    ".cpp": "text",
    ".sh": "text",
}

_MEDIA_TYPES = {
    "markdown": "text/markdown",
    "csv": "text/csv",
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "notebook": "application/x-ipynb+json",
    "image": "application/octet-stream",
    "archive": "application/zip",
    "pdb": "chemical/x-pdb",
    "model3d": "application/octet-stream",
    "text": "text/plain",
    "unknown": "application/octet-stream",
}

_IMAGE_MEDIA = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".webp": "image/webp",
    ".bmp": "image/bmp",
    ".gif": "image/gif",
}


def classify_preview(path: Path) -> str:
    return _SUFFIX_FAMILY.get(path.suffix.lower(), "unknown")


def _result(
    *,
    family: str,
    title: str,
    text: str = "",
    truncated: bool = False,
    metadata: dict[str, Any] | None = None,
    error: str | None = None,
    media_type: str | None = None,
    path: Path | None = None,
) -> dict[str, Any]:
    if media_type is None:
        if family == "image" and path is not None:
            media_type = _IMAGE_MEDIA.get(path.suffix.lower(), "application/octet-stream")
        else:
            media_type = _MEDIA_TYPES.get(family, "application/octet-stream")
    return {
        "family": family,
        "title": title,
        "truncated": truncated,
        "text": text,
        "metadata": metadata or {},
        "error": error,
        "media_type": media_type,
    }


def build_preview(path: Path, *, max_bytes: int = _DEFAULT_MAX_BYTES) -> dict[str, Any]:
    path = Path(path)
    family = classify_preview(path)
    title = path.name
    try:
        if family == "notebook":
            return _preview_notebook(path)
        if family == "pdb":
            return _preview_pdb(path, max_bytes=max_bytes)
        if family == "model3d":
            return _preview_3d(path)
        if family == "archive":
            return _preview_archive(path)
        if family == "pdf":
            return _preview_pdf(path)
        if family == "docx":
            return _preview_docx(path)
        if family == "xlsx":
            return _preview_xlsx(path)
        if family == "pptx":
            return _preview_pptx(path)
        if family == "image":
            return _preview_image(path)
        if family == "csv":
            return _preview_csv(path)
        return _preview_text(path, family=family, max_bytes=max_bytes)
    except Exception as exc:
        return _result(
            family=family,
            title=title,
            error=str(exc),
            path=path,
        )


def _preview_text(path: Path, *, family: str, max_bytes: int) -> dict[str, Any]:
    raw = path.read_bytes()
    truncated = len(raw) > max_bytes
    raw = raw[:max_bytes]
    text = raw.decode("utf-8", errors="replace")
    if truncated:
        text += "\n\n… preview truncated; open the file for the complete contents.\n"
    return _result(family=family if family != "unknown" else "text", title=path.name, text=text, truncated=truncated, path=path)


def _preview_notebook(path: Path) -> dict[str, Any]:
    nb = json.loads(path.read_text(encoding="utf-8", errors="replace"))
    cells = nb.get("cells", [])
    lines = [f"📓 Jupyter Notebook · {len(cells)} cell(s)", "─" * 45, ""]
    shown = cells[:20]
    for idx, cell in enumerate(shown, 1):
        cell_type = cell.get("cell_type", "code")
        source = "".join(cell.get("source", []))
        if cell_type == "markdown":
            lines.append(f"─── [Cell {idx}: Markdown] ───")
            lines.append(source)
            lines.append("")
        else:
            lines.append(f"─── [Cell {idx}: Code] ───")
            lines.append(source)
            for out in cell.get("outputs", [])[:3]:
                out_type = out.get("output_type", "")
                if out_type in ("stream", "execute_result"):
                    text_out = "".join(out.get("text", []) or out.get("data", {}).get("text/plain", []))
                    if text_out.strip():
                        lines.append("Out: " + text_out.strip()[:400])
            lines.append("")
    truncated = len(cells) > 20
    if truncated:
        lines.append(f"… {len(cells) - 20} more cell(s) omitted.")
    return _result(
        family="notebook",
        title=f"Jupyter Notebook · {len(cells)} cell(s)",
        text="\n".join(lines),
        truncated=truncated,
        metadata={"cells": len(cells)},
        path=path,
    )


def _preview_pdb(path: Path, *, max_bytes: int) -> dict[str, Any]:
    with path.open("r", encoding="utf-8", errors="replace") as handle:
        lines = handle.readlines(max_bytes)
    header = ""
    title = ""
    expdta = ""
    resolution = ""
    chains: set[str] = set()
    atom_count = 0
    hetatm_count = 0
    b_factors: list[float] = []
    for line in lines:
        rec = line[:6].strip()
        if rec == "HEADER":
            header = line[10:50].strip()
        elif rec == "TITLE":
            title += line[10:70].strip() + " "
        elif rec == "EXPDTA":
            expdta = line[10:70].strip()
        elif rec == "REMARK" and "RESOLUTION." in line:
            resolution = line.strip()
        elif rec in ("ATOM", "HETATM"):
            if rec == "ATOM":
                atom_count += 1
            else:
                hetatm_count += 1
            chain_id = line[21:22].strip()
            if chain_id:
                chains.add(chain_id)
            try:
                b_factors.append(float(line[60:66].strip()))
            except ValueError:
                pass
    parts = [f"🧬 PDB Biological Structure · {path.name}", "─" * 45]
    if header:
        parts.append(f"Header: {header}")
    if title:
        parts.append(f"Title: {title.strip()}")
    if expdta:
        parts.append(f"Method: {expdta}")
    if resolution:
        parts.append(resolution)
    parts.append(f"Chains: {', '.join(sorted(chains)) if chains else 'N/A'}")
    parts.append(f"Atoms: {atom_count:,} (ATOM) | {hetatm_count:,} (HETATM)")
    avg_b = None
    if b_factors:
        avg_b = sum(b_factors) / len(b_factors)
        parts.append(f"Avg B-Factor / pLDDT: {avg_b:.2f}")
    return _result(
        family="pdb",
        title=f"PDB Biological Structure · {path.name}",
        text="\n".join(parts),
        metadata={"atoms": atom_count, "hetatm": hetatm_count, "avg_b": avg_b},
        path=path,
    )


def _preview_3d(path: Path) -> dict[str, Any]:
    suffix = path.suffix.lower()
    parts = [f"🧊 3D Asset · {path.name} ({suffix.upper()})", "─" * 45]
    meta: dict[str, Any] = {"suffix": suffix}
    if suffix == ".obj":
        content = path.read_text(encoding="utf-8", errors="replace")
        v_count = sum(1 for line in content.splitlines() if line.startswith("v "))
        f_count = sum(1 for line in content.splitlines() if line.startswith("f "))
        parts.append(f"Format: Wavefront OBJ\nVertices: {v_count:,}\nFaces: {f_count:,}")
        meta.update({"vertices": v_count, "faces": f_count})
    elif suffix == ".stl":
        header = path.read_bytes()[:80]
        is_ascii = header.startswith(b"solid")
        if is_ascii:
            facets = sum(1 for line in path.read_text(encoding="utf-8", errors="replace").splitlines() if "facet normal" in line)
            parts.append(f"Format: ASCII STL\nFacets: {facets:,}")
            meta["facets"] = facets
        else:
            import struct

            with path.open("rb") as handle:
                handle.seek(80)
                num_triangles = struct.unpack("<I", handle.read(4))[0]
            parts.append(f"Format: Binary STL\nTriangles: {num_triangles:,}")
            meta["triangles"] = num_triangles
    elif suffix == ".gltf":
        data = json.loads(path.read_text(encoding="utf-8", errors="replace"))
        asset = data.get("asset", {})
        parts.append(f"Format: glTF {asset.get('version', '2.0')}")
        parts.append(f"Meshes: {len(data.get('meshes', []))}\nMaterials: {len(data.get('materials', []))}\nNodes: {len(data.get('nodes', []))}")
    elif suffix == ".glb":
        parts.append(f"Format: Binary glTF (GLB)\nSize: {path.stat().st_size:,} bytes")
    return _result(family="model3d", title=f"3D Asset · {path.name}", text="\n".join(parts), metadata=meta, path=path)


def _preview_archive(path: Path) -> dict[str, Any]:
    suffix = path.suffix.lower()
    parts = [f"📦 Archive · {path.name}", "─" * 45]
    if suffix == ".zip":
        import zipfile

        with zipfile.ZipFile(path, "r") as archive:
            infos = archive.infolist()
            parts.append(f"Entries: {len(infos)}\n")
            for info in infos[:25]:
                size_str = f"{info.file_size:,} B" if not info.is_dir() else "<DIR>"
                parts.append(f"{info.filename:<40} {size_str:>12}")
            if len(infos) > 25:
                parts.append(f"\n… {len(infos) - 25} more entries omitted.")
            names = [info.filename for info in infos]
    else:
        import tarfile

        mode = "r:gz" if suffix in {".gz", ".tgz"} else "r"
        with tarfile.open(path, mode) as archive:
            members = archive.getmembers()
            parts.append(f"Entries: {len(members)}\n")
            for member in members[:25]:
                size_str = f"{member.size:,} B" if member.isfile() else "<DIR>"
                parts.append(f"{member.name:<40} {size_str:>12}")
            if len(members) > 25:
                parts.append(f"\n… {len(members) - 25} more entries omitted.")
            names = [member.name for member in members]
    return _result(
        family="archive",
        title=f"Archive · {path.name}",
        text="\n".join(parts),
        metadata={"entries": names},
        path=path,
    )


def _preview_pdf(path: Path) -> dict[str, Any]:
    try:
        import pypdf
    except ImportError:
        return _result(
            family="pdf",
            title="PDF Document (.pdf)",
            text="Install optional package 'pypdf' for full content preview:\n   pip install frontier-agent[document-readers]",
            error="missing pypdf",
            path=path,
        )
    try:
        reader = pypdf.PdfReader(str(path))
        num_pages = len(reader.pages)
        meta = reader.metadata or {}
        title = meta.get("/Title", "") or path.name
        author = meta.get("/Author", "")
        parts = [f"📄 PDF Document · {num_pages} page(s)"]
        if title and title != path.name:
            parts.append(f"Title: {title}")
        if author:
            parts.append(f"Author: {author}")
        parts.append("─" * 45)
        max_pages = min(num_pages, 5)
        for i in range(max_pages):
            parts.append(f"--- Page {i + 1} ---")
            page_text = reader.pages[i].extract_text() or ""
            parts.append(page_text.strip()[:2000] if page_text.strip() else "[No extractable text on this page]")
        truncated = num_pages > max_pages
        if truncated:
            parts.append(f"… {num_pages - max_pages} more page(s) not shown.")
        return _result(
            family="pdf",
            title="PDF Document",
            text="\n".join(parts),
            truncated=truncated,
            metadata={"pages": num_pages, "title": title, "author": author},
            path=path,
        )
    except Exception as exc:
        return _result(family="pdf", title="PDF Document", error=f"Could not read PDF: {exc}", path=path)


def _preview_docx(path: Path) -> dict[str, Any]:
    try:
        import docx
    except ImportError:
        return _result(
            family="docx",
            title="Word Document (.docx)",
            text="Install optional package 'python-docx' for full content preview:\n   pip install frontier-agent[document-readers]",
            error="missing python-docx",
            path=path,
        )
    try:
        doc = docx.Document(str(path))
        lines = [f"# {path.name}", f"*Word Document ({len(doc.paragraphs)} paragraphs, {len(doc.tables)} tables)*", "---"]
        for paragraph in doc.paragraphs:
            if not paragraph.text.strip():
                continue
            style_name = (paragraph.style.name or "").lower() if paragraph.style else ""
            if "heading 1" in style_name:
                lines.append(f"# {paragraph.text}")
            elif "heading 2" in style_name:
                lines.append(f"## {paragraph.text}")
            elif "heading 3" in style_name:
                lines.append(f"### {paragraph.text}")
            else:
                lines.append(paragraph.text)
        return _result(
            family="docx",
            title="Word Document",
            text="\n\n".join(lines),
            metadata={"paragraphs": len(doc.paragraphs), "tables": len(doc.tables)},
            path=path,
        )
    except Exception as exc:
        return _result(family="docx", title="Word Document", error=f"Could not read Word document: {exc}", path=path)


def _preview_xlsx(path: Path) -> dict[str, Any]:
    try:
        import openpyxl
    except ImportError:
        return _result(
            family="xlsx",
            title="Excel Spreadsheet (.xlsx)",
            text="Install optional package 'openpyxl' for full content preview:\n   pip install frontier-agent[sandbox]",
            error="missing openpyxl",
            path=path,
        )
    try:
        workbook = openpyxl.load_workbook(filename=str(path), read_only=True, data_only=True)
        sheet_names = workbook.sheetnames
        sheet_name = sheet_names[0] if sheet_names else "Sheet1"
        worksheet = workbook[sheet_name]
        rows = [[("" if cell is None else str(cell)) for cell in row] for row in worksheet.iter_rows(values_only=True)]
        workbook.close()
        shown = rows[:26]
        text_rows = ["\t".join(row[:8]) for row in shown]
        return _result(
            family="xlsx",
            title=f"Excel: {path.name}",
            text="\n".join(text_rows),
            metadata={"sheets": sheet_names, "rows": shown},
            path=path,
        )
    except Exception as exc:
        return _result(family="xlsx", title="Excel Spreadsheet", error=f"Could not read Excel file: {exc}", path=path)


def _preview_pptx(path: Path) -> dict[str, Any]:
    try:
        import pptx
    except ImportError:
        return _result(
            family="pptx",
            title="PowerPoint Presentation (.pptx)",
            text="Install optional package 'python-pptx' for full content preview:\n   pip install frontier-agent[document-readers]",
            error="missing python-pptx",
            path=path,
        )
    try:
        presentation = pptx.Presentation(str(path))
        lines = [f"# {path.name}", f"*PowerPoint Presentation ({len(presentation.slides)} slides)*", "---"]
        for idx, slide in enumerate(presentation.slides, 1):
            slide_title = "Untitled Slide"
            if slide.shapes.title and slide.shapes.title.text:
                slide_title = slide.shapes.title.text.strip()
            lines.append(f"### Slide {idx}: {slide_title}")
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    for paragraph in shape.text_frame.paragraphs:  # pyright: ignore[reportAttributeAccessIssue]
                        if paragraph.text.strip():
                            lines.append(f"- {paragraph.text.strip()}")
        return _result(
            family="pptx",
            title="PowerPoint Presentation",
            text="\n".join(lines),
            metadata={"slides": len(presentation.slides)},
            path=path,
        )
    except Exception as exc:
        return _result(family="pptx", title="PowerPoint Presentation", error=f"Could not read PowerPoint presentation: {exc}", path=path)


def _preview_image(path: Path) -> dict[str, Any]:
    try:
        from PIL import Image
    except ImportError:
        return _result(
            family="image",
            title=f"Image · {path.name}",
            text="Install optional package 'Pillow' for image preview:\n   pip install Pillow",
            error="missing Pillow",
            path=path,
        )
    try:
        with Image.open(path) as image:
            width, height = image.size
            mode = image.mode
        return _result(
            family="image",
            title=f"Image · {width}x{height} ({mode}) · {path.name}",
            text=f"{width}x{height} {mode}",
            metadata={"width": width, "height": height, "mode": mode},
            path=path,
        )
    except Exception as exc:
        return _result(family="image", title=f"Image · {path.name}", error=f"Could not load image: {exc}", path=path)


def _preview_csv(path: Path) -> dict[str, Any]:
    import csv

    with path.open("r", encoding="utf-8", errors="replace") as handle:
        rows = list(csv.reader(handle))
    if not rows:
        return _result(family="csv", title=path.name, text="Empty CSV file", path=path)
    shown = [row[:10] for row in rows[:31]]
    text = "\n".join("\t".join(str(cell) for cell in row) for row in shown)
    return _result(family="csv", title=path.name, text=text, metadata={"rows": shown}, path=path)
